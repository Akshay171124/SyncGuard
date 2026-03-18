"""Generate SyncGuard project presentation — 7 slides."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BLUE = RGBColor(0x1A, 0x52, 0x76)
MEDIUM_BLUE = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BLUE = RGBColor(0x34, 0x98, 0xDB)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0x95, 0xA5, 0xA6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
BG_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)

TOTAL_SLIDES = 7


def set_slide_bg(slide, color=BG_LIGHT):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title_text, subtitle_text=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, Inches(1.1),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), Inches(12), Inches(0.6))
    p = txBox.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.68), Inches(12), Inches(0.35))
        p2 = txBox2.text_frame.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)


def add_slide_number(slide, num):
    txBox = slide.shapes.add_textbox(Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = f"{num}/{TOTAL_SLIDES}"
    p.font.size = Pt(10)
    p.font.color.rgb = MED_GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_content_box(slide, left, top, width, height, title, items, box_color=MEDIUM_BLUE, font_size=11):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
    shape.line.color.rgb = box_color
    shape.line.width = Pt(2)

    txBox = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.08),
        Inches(width - 0.3), Inches(0.35),
    )
    p = txBox.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = box_color

    txBox2 = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.42),
        Inches(width - 0.4), Inches(height - 0.5),
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(2)


def add_body_text(slide, left, top, width, height, texts, font_size=14):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(texts):
        if isinstance(item, str):
            text, bold, color, size = item, False, DARK_GRAY, font_size
        elif len(item) == 2:
            text, bold = item
            color, size = DARK_GRAY, font_size
        elif len(item) == 3:
            text, bold, color = item
            size = font_size
        else:
            text, bold, color, size = item

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(4)


# ========== SLIDE 1: Title ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BLUE)

txBox = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "SyncGuard"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Contrastive Audio-Visual Deepfake Detection\nvia Temporal Phoneme-Face Coherence"
p2.font.size = Pt(24)
p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(16)

shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.3), Inches(5.333), Inches(0.03),
)
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_BLUE
shape.line.fill.background()

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.7), Inches(11), Inches(2.2))
tf2 = txBox2.text_frame

p3 = tf2.paragraphs[0]
p3.text = "CS 5330 — Computer Vision  |  Northeastern University, Khoury College"
p3.font.size = Pt(16)
p3.font.color.rgb = RGBColor(0x99, 0xBB, 0xDD)
p3.alignment = PP_ALIGN.CENTER

p4 = tf2.add_paragraph()
p4.text = "Akshay Prajapati  |  Ritik  |  Atharva"
p4.font.size = Pt(16)
p4.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
p4.alignment = PP_ALIGN.CENTER
p4.space_before = Pt(12)

p5 = tf2.add_paragraph()
p5.text = "March 2026"
p5.font.size = Pt(13)
p5.font.color.rgb = RGBColor(0x77, 0x99, 0xBB)
p5.alignment = PP_ALIGN.CENTER
p5.space_before = Pt(10)


# ========== SLIDE 2: Problem & Motivation ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, "The Problem", "Why current deepfake detection falls short")
add_slide_number(slide, 2)

# Left column: the problem
add_body_text(slide, 0.6, 1.3, 5.8, 5.5, [
    ("Deepfakes are evolving faster than detectors", True, DARK_BLUE, 18),
    "",
    ("Current approaches rely on visual artifacts:", True, DARK_GRAY, 14),
    "  Blending boundaries, texture inconsistencies, flickering",
    "  These artifacts are generator-specific and disappear as",
    "  synthesis improves — detectors trained on one method",
    "  fail on the next",
    "",
    ("What's missing:", True, RED, 14),
    "  Most detectors ignore audio entirely",
    "  Even AV methods check global consistency, not",
    "  frame-level temporal synchronization",
    "",
    ("Our approach: exploit the physics of speech", True, GREEN, 14),
    "  Real speech → tight temporal coupling between",
    "  lip motion and audio phonemes",
    "  Fakes optimize appearance, not fine-grained sync",
    "  This is a generator-agnostic signal",
], font_size=13)

# Right column: visual comparison
add_content_box(slide, 6.8, 1.3, 5.8, 2.4,
    "Generator-Specific Detection (Fragile)",
    [
        "Trained on FaceSwap → fails on Wav2Lip",
        "Trained on Wav2Lip → fails on FaceSwap",
        "Visual artifacts change with every new generator",
        "Cross-dataset AUC often drops 20-30%",
        "",
        "Detectors are playing catch-up with generators",
    ], RED, font_size=12)

add_content_box(slide, 6.8, 3.95, 5.8, 3.0,
    "SyncGuard: Generator-Agnostic (Robust)",
    [
        "Key insight: speech production is governed by",
        "biomechanics — not generator architecture",
        "",
        "Lip rounding for /o/ always co-occurs with the",
        "corresponding acoustic formant in real speech",
        "",
        "No generator perfectly replicates this coupling",
        "at the frame level (20-50ms resolution)",
        "",
        "Goal: AUC ≥ 0.88 in-domain, ≥ 0.72 zero-shot",
    ], GREEN, font_size=12)


# ========== SLIDE 3: How SyncGuard Works ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, "How SyncGuard Works", "Two-stream contrastive learning for AV synchrony detection")
add_slide_number(slide, 3)

# Architecture diagram
diagram_path = "outputs/visualizations/architecture_diagram.png"
if os.path.exists(diagram_path):
    slide.shapes.add_picture(
        diagram_path,
        Inches(0.2), Inches(1.2),
        Inches(12.9), Inches(6.1),
    )


# ========== SLIDE 4: Technical Deep Dive ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, "Technical Deep Dive", "Model components, training strategy, and design choices")
add_slide_number(slide, 4)

# Row 1: encoders + classifier
add_content_box(slide, 0.3, 1.3, 3.1, 3.0,
    "Visual Stream",
    [
        "AV-HuBERT (pretrained lip-reading)",
        "  Trained on 1,759 hrs of speech",
        "  3D Conv (5×7×7) + ResNet-18",
        "",
        "Projection: Linear→ReLU→Linear",
        "Output: v_t ∈ R^256, L2-normalized",
        "",
        "Why: lip-reading pretraining encodes",
        "  exactly the visual phonetics we need",
    ], MEDIUM_BLUE, font_size=11)

add_content_box(slide, 3.55, 1.3, 3.1, 3.0,
    "Audio Stream",
    [
        "Wav2Vec 2.0 (facebook, frozen)",
        "  Layer 9 hidden states extracted",
        "  (best phonetic encoding per",
        "  Pasad et al. 2021)",
        "",
        "Projection: Linear→ReLU→Linear",
        "Output: a_t ∈ R^256, L2-normalized",
        "",
        "Why frozen: 94M params saved from",
        "  gradient computation → fast training",
    ], LIGHT_BLUE, font_size=11)

add_content_box(slide, 6.8, 1.3, 3.1, 3.0,
    "Sync Score + Classifier",
    [
        "Frame-level cosine similarity:",
        "  s(t) = cos(v_t, a_t) ∈ [-1, 1]",
        "",
        "Real: s(t) ≈ 0.7-0.9 (synchronized)",
        "Fake: s(t) ≈ 0.1-0.4 (desynchronized)",
        "",
        "Bi-LSTM (2-layer, hidden=128)",
        "  Captures temporal patterns in s(t)",
        "  Mean+Max pool → binary decision",
    ], PURPLE, font_size=11)

add_content_box(slide, 10.05, 1.3, 3.0, 3.0,
    "Training Strategy",
    [
        "Phase 1: Contrastive Pretrain",
        "  Real data only, InfoNCE loss",
        "  MoCo queue (4096 negatives)",
        "  Learn AV alignment, 20 epochs",
        "",
        "Phase 2: Fine-tune",
        "  FakeAVCeleb (real + fake)",
        "  L_nce + 0.5·L_temp + L_cls",
        "  Hard neg mining 0%→20%",
    ], GREEN, font_size=11)

# Row 2: losses + key decisions
add_content_box(slide, 0.3, 4.5, 4.2, 2.6,
    "Loss Design",
    [
        "InfoNCE: frame-level contrastive, v_t vs MoCo negatives",
        "  Learnable temperature τ (init=0.07, clamped [0.01, 0.5])",
        "",
        "Temporal Consistency (real clips only):",
        "  L_temp = Σ_t ||(Δv_t) - (Δa_t)||²",
        "  Penalizes divergent rate-of-change in embeddings",
        "",
        "BCE: standard classification loss on logits",
    ], ORANGE, font_size=11)

add_content_box(slide, 4.7, 4.5, 4.2, 2.6,
    "Ablation Dimensions",
    [
        "Visual encoder:  AV-HuBERT  |  ResNet-18  |  SyncNet",
        "  (Does lip-reading pretraining matter?)",
        "",
        "Wav2Vec layer:  3  |  5  |  7  |  9  |  11",
        "  (Which layer best encodes phonetics?)",
        "",
        "Classifier:  Bi-LSTM  |  1D-CNN  |  Statistical",
        "  (Does temporal modeling matter?)",
    ], PURPLE, font_size=11)

add_content_box(slide, 9.1, 4.5, 4.1, 2.6,
    "Tools & Infrastructure",
    [
        "PyTorch, HuggingFace, fairseq",
        "OpenCV, RetinaFace, MediaPipe",
        "Silero-VAD, torchaudio",
        "",
        "Northeastern HPC Explorer",
        "  NVIDIA H200 (140GB) GPUs",
        "  SLURM scheduler",
        "YAML config-driven pipeline",
    ], DARK_BLUE, font_size=11)


# ========== SLIDE 5: Data Pipeline ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, "Data & Preprocessing", "Datasets, preprocessing pipeline, and evaluation strategy")
add_slide_number(slide, 5)

# Datasets - visual layout
datasets = [
    ("FakeAVCeleb", "Primary Dataset", "19,500 clips\n4 manipulation types\nSpeaker-disjoint splits", MEDIUM_BLUE, "Train / Val / Test"),
    ("VoxCeleb2", "Pretraining Data", "~500 hrs (subset)\nReal speech only\nDiverse speakers", GREEN, "Phase 1 Pretrain"),
    ("CelebDF-v2", "Zero-Shot Test", "6,229 clips\nDifferent generator\nNo training exposure", ORANGE, "Generalization"),
    ("DFDC", "Zero-Shot Test", "~5K test clips\nFacebook challenge\nHardest benchmark", PURPLE, "Generalization"),
    ("Wav2Lip", "Adversarial Test", "~500 clips\nSync-optimized fakes\nSelf-generated", RED, "Robustness"),
]

x_start = 0.3
for i, (name, subtitle, desc, color, role) in enumerate(datasets):
    x = x_start + i * 2.6
    # Header
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(1.3), Inches(2.4), Inches(0.7),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Body
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(2.0), Inches(2.4), Inches(1.7),
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFB)
    shape2.line.color.rgb = color
    shape2.line.width = Pt(1.5)

    txBox = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(2.05), Inches(2.1), Inches(1.5),
    )
    tf2 = txBox.text_frame
    tf2.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        pp = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
        pp.text = line
        pp.font.size = Pt(11)
        pp.font.color.rgb = DARK_GRAY
        pp.alignment = PP_ALIGN.CENTER
        pp.space_after = Pt(2)

    # Role tag
    txBox3 = slide.shapes.add_textbox(
        Inches(x), Inches(3.6), Inches(2.4), Inches(0.35),
    )
    p3 = txBox3.text_frame.paragraphs[0]
    p3.text = role
    p3.font.size = Pt(10)
    p3.font.bold = True
    p3.font.color.rgb = color
    p3.alignment = PP_ALIGN.CENTER

# Preprocessing pipeline
add_content_box(slide, 0.3, 4.2, 6.6, 3.0,
    "Preprocessing Pipeline",
    [
        "1. Video → RetinaFace face detection (confidence > 0.8)",
        "2. MediaPipe FaceMesh → 468 landmarks → mouth ROI crop",
        "3. Resize to 96×96 grayscale, normalize [0, 1]",
        "4. Audio → FFmpeg extract → resample to 16kHz mono",
        "5. Silero-VAD → speech activity detection mask",
        "6. Temporal alignment: 25fps visual ↔ 49Hz Wav2Vec rate",
        "",
        "Output per clip: mouth_crops.npy  |  audio.wav  |  speech_mask.npy",
    ], GREEN, font_size=11)

# FakeAVCeleb breakdown
add_content_box(slide, 7.1, 4.2, 5.9, 3.0,
    "FakeAVCeleb — 4 Manipulation Categories",
    [
        "RV-RA  Real Video + Real Audio      (genuine)",
        "FV-RA  Fake Video + Real Audio      (face swap)",
        "RV-FA  Real Video + Fake Audio      (voice clone)",
        "FV-FA  Fake Video + Fake Audio      (full deepfake)",
        "",
        "Evaluation strategy: per-category AUC breakdown",
        "  Hypothesis: FV-FA easiest (both streams off),",
        "  RV-FA hardest (only audio is fake — subtle desync)",
    ], MEDIUM_BLUE, font_size=11)


# ========== SLIDE 6: Current Progress ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, "Where We Are", "Implementation progress and early verification")
add_slide_number(slide, 6)

# Left: what's built
add_body_text(slide, 0.5, 1.2, 5.5, 0.5, [
    ("Implemented & Tested", True, DARK_BLUE, 18),
])

components = [
    "Preprocessing pipeline (RetinaFace + MediaPipe + VAD)",
    "Visual encoder — AV-HuBERT + 2 ablation variants",
    "Audio encoder — Wav2Vec 2.0, frozen backbone, layer 9",
    "Temporal classifier — Bi-LSTM + 2 ablation variants",
    "Full model integration — 107M params, 13M trainable",
    "Loss functions — InfoNCE + MoCo, temporal consistency, BCE",
    "Training dataset — speaker-disjoint splits, hard neg mining",
    "Phase 1 pretraining loop — cosine LR, checkpointing",
    "Phase 2 fine-tuning loop — combined loss, early stopping",
    "CLI scripts — ready to launch on HPC",
]
y = 1.75
for comp in components:
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(5.5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"  {comp}"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY

    # Checkmark
    txBox2 = slide.shapes.add_textbox(Inches(0.4), Inches(y), Inches(0.3), Inches(0.3))
    p2 = txBox2.text_frame.paragraphs[0]
    p2.text = "✓"
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = GREEN

    y += 0.33

# Right: verification results
add_content_box(slide, 6.6, 1.2, 6.3, 2.3,
    "Shape & Gradient Verification",
    [
        "End-to-end: (B,T,1,96,96) → (B,T,256) → (B,T) → (B,1) ✓",
        "L2 normalization on both embedding streams ✓",
        "Gradient flow through all trainable projection heads ✓",
        "Frozen Wav2Vec backbone: 94M params, zero gradients ✓",
        "Variable-length batches with padding + boolean masks ✓",
    ], MEDIUM_BLUE, font_size=11)

add_content_box(slide, 6.6, 3.7, 6.3, 1.8,
    "Training Loop Tests (CPU, synthetic data)",
    [
        "Pretraining: InfoNCE loss 5.02 → 4.87 (decreasing ✓)",
        "Sync-score: -0.03 → 0.18 (learning alignment ✓)",
        "Fine-tuning: all 3 loss components compute correctly ✓",
        "Checkpoint save/load with full state (model + optimizer + scheduler) ✓",
    ], GREEN, font_size=11)

add_content_box(slide, 6.6, 5.7, 6.3, 1.5,
    "Critical Bug Caught Early",
    [
        "Wav2Vec 2.0 produces NaN on zero-padded waveforms in train mode",
        "Cause: group normalization computes stats over zero regions",
        "Fix: force frozen backbone to eval() during forward pass",
        "Caught via systematic CPU testing — would have silently failed on HPC",
    ], RED, font_size=11)


# ========== SLIDE 7: What's Next ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, "Road Ahead", "From implementation to results — April 13 deadline")
add_slide_number(slide, 7)

# Timeline
phases = [
    ("This Week", "Data + Eval Setup",
     "Transfer datasets to HPC, preprocess FakeAVCeleb, build evaluation framework (AUC, EER, pAUC)",
     MEDIUM_BLUE),
    ("Mar 16–21", "Contrastive Pretraining",
     "Phase 1: train encoders on real speech data, 20 epochs on H200 GPU. Target: avg sync-score > 0.5",
     GREEN),
    ("Mar 22–28", "Fine-tuning + Evaluation",
     "Phase 2: fine-tune on FakeAVCeleb, evaluate on CelebDF-v2 & DFDC. Target: AUC ≥ 0.88 / ≥ 0.72",
     ORANGE),
    ("Mar 29–Apr 5", "Ablation Study",
     "8 experiments: visual encoder (×3), Wav2Vec layer (×5), classifier (×3), hard negatives",
     PURPLE),
    ("Apr 6–13", "Final Deliverables",
     "Paper, poster, video demo, sync-score visualizations",
     RED),
]

y = 1.3
for date, title, desc, color in phases:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(y), Inches(2.0), Inches(0.5),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = date
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    txBox = slide.shapes.add_textbox(Inches(2.7), Inches(y + 0.0), Inches(2.5), Inches(0.5))
    p2 = txBox.text_frame.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(15)
    p2.font.bold = True
    p2.font.color.rgb = DARK_GRAY
    txBox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    txBox2 = slide.shapes.add_textbox(Inches(5.3), Inches(y + 0.0), Inches(7.5), Inches(0.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    p3.text = desc
    p3.font.size = Pt(12)
    p3.font.color.rgb = DARK_GRAY
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE

    if date != "Apr 6–13":
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.48), Inches(y + 0.5), Inches(0.04), Inches(0.5),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        line.line.fill.background()

    y += 1.0

# Bottom: expected outcomes + thank you
add_content_box(slide, 0.3, 6.5, 5.5, 0.8,
    "Expected Outcomes",
    [
        "FakeAVCeleb AUC ≥ 0.88  |  CelebDF ≥ 0.79  |  DFDC ≥ 0.72  |  Wav2Lip: report",
    ], DARK_BLUE, font_size=12)

txBox = slide.shapes.add_textbox(Inches(7.0), Inches(6.4), Inches(6.0), Inches(1.0))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Thank You — Questions?"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "github.com/Akshay171124/SyncGuard"
p2.font.size = Pt(11)
p2.font.color.rgb = MED_GRAY
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(4)


# Save
output_path = "outputs/SyncGuard_Phase3_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {TOTAL_SLIDES}")
